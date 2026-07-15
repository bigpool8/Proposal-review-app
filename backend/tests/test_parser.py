import io
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches
from docx import Document

from app.services.parser import parse_file

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PDF_MIME = "application/pdf"


# ── 테스트 파일 생성 헬퍼 ──────────────────────────────────


def make_pptx(texts_per_slide: list[list[str]]) -> str:
    """슬라이드별 텍스트 목록으로 임시 PPTX 파일 생성."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout
    for texts in texts_per_slide:
        slide = prs.slides.add_slide(blank_layout)
        for text in texts:
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tb.text_frame.text = text
    f = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(f.name)
    f.close()
    return f.name


def make_docx(paragraphs: list[str], table_cells: list[str] | None = None) -> str:
    """단락 목록으로 임시 DOCX 파일 생성. table_cells가 있으면 표도 추가."""
    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    if table_cells:
        table = doc.add_table(rows=1, cols=len(table_cells))
        for i, text in enumerate(table_cells):
            table.rows[0].cells[i].text = text
    f = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
    doc.save(f.name)
    f.close()
    return f.name


def make_test_image_bytes(size: tuple[int, int] = (60, 60)) -> bytes:
    """MIN_IMAGE_BYTES(1KB) 기준을 넘기는 임의 패턴의 PNG 이미지 바이트 생성."""
    from PIL import Image

    img = Image.effect_noise(size, 50).convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def make_pdf(pages_text: list[str]) -> str:
    """페이지별 텍스트(ASCII/Latin-1)로 임시 PDF 파일 생성 (fpdf2 사용)."""
    from fpdf import FPDF

    pdf = FPDF()
    for text in pages_text:
        pdf.add_page()
        if text:
            pdf.set_font("Helvetica", size=12)
            pdf.cell(0, 10, text=text)
    f = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    pdf.output(f.name)
    f.close()
    return f.name


# ── PPTX 테스트 ──────────────────────────────────────────


class TestPptxParser:
    def test_slide_count(self):
        path = make_pptx([["슬라이드 1"], ["슬라이드 2"], ["슬라이드 3"]])
        try:
            r = parse_file(path, PPTX_MIME, "f1", "test.pptx", "qualitative")
            assert r["parse_error"] is None
            assert r["total_pages"] == 3
            assert len(r["pages"]) == 3
        finally:
            os.unlink(path)

    def test_page_numbers_one_based(self):
        path = make_pptx([["첫째"], ["둘째"]])
        try:
            r = parse_file(path, PPTX_MIME, "f2", "test.pptx", "qualitative")
            assert r["pages"][0]["page_number"] == 1
            assert r["pages"][1]["page_number"] == 2
        finally:
            os.unlink(path)

    def test_text_extracted(self):
        path = make_pptx([["안녕하세요 테스트"]])
        try:
            r = parse_file(path, PPTX_MIME, "f3", "test.pptx", "presentation")
            assert "안녕하세요 테스트" in r["pages"][0]["text"]
        finally:
            os.unlink(path)

    def test_empty_slide_included(self):
        path = make_pptx([["텍스트 있음"], []])
        try:
            r = parse_file(path, PPTX_MIME, "f4", "test.pptx", "qualitative")
            assert r["total_pages"] == 2
            assert r["pages"][1]["text"] == ""
        finally:
            os.unlink(path)

    def test_result_metadata(self):
        path = make_pptx([["슬라이드"]])
        try:
            r = parse_file(path, PPTX_MIME, "fid-meta", "slides.pptx", "quantitative")
            assert r["file_id"] == "fid-meta"
            assert r["original_filename"] == "slides.pptx"
            assert r["proposal_type"] == "quantitative"
        finally:
            os.unlink(path)


class TestPptxImageOcclusion:
    """불투명 도형에 가려진 그림은 그 영역만 흰색으로 마스킹된 뒤 추출되어야 한다."""

    def _pixel(self, image_bytes: bytes, xy: tuple[int, int]):
        from PIL import Image
        return Image.open(io.BytesIO(image_bytes)).convert("RGB").getpixel(xy)

    def test_fully_covered_by_solid_shape_masked_white(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(
            io.BytesIO(make_test_image_bytes()), Inches(1), Inches(1), Inches(2), Inches(2)
        )
        cover = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, pic.left, pic.top, pic.width, pic.height
        )
        cover.fill.solid()
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        prs.save(path)
        try:
            r = parse_file(path, PPTX_MIME, "occ1", "test.pptx", "qualitative")
            assert len(r["images"]) == 1
            assert self._pixel(r["images"][0]["image_bytes"], (30, 30)) == (255, 255, 255)
        finally:
            os.unlink(path)

    def test_partially_covered_masks_only_that_region(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        original_bytes = make_test_image_bytes((60, 60))
        pic = slide.shapes.add_picture(
            io.BytesIO(original_bytes), Inches(1), Inches(1), Inches(2), Inches(2)
        )
        # 왼쪽 절반만 덮음
        cover = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, pic.left, pic.top, Emu(int(pic.width / 2)), pic.height
        )
        cover.fill.solid()
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        prs.save(path)
        try:
            r = parse_file(path, PPTX_MIME, "occ2", "test.pptx", "qualitative")
            assert len(r["images"]) == 1
            masked_bytes = r["images"][0]["image_bytes"]
            # 덮인 왼쪽은 흰색으로 마스킹됨
            assert self._pixel(masked_bytes, (10, 30)) == (255, 255, 255)
            # 덮이지 않은 오른쪽은 원본 픽셀 그대로 유지됨
            assert self._pixel(masked_bytes, (50, 30)) == self._pixel(original_bytes, (50, 30))
        finally:
            os.unlink(path)

    def test_covered_by_another_picture_not_masked(self):
        # 그림 위에 다른 그림(액자처럼 가운데가 뚫려 있을 수 있음)이 겹쳐도
        # 픽셀 투명도를 확인할 수 없으므로 안전하게 마스킹하지 않아야 한다.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        original_bytes = make_test_image_bytes()
        pic = slide.shapes.add_picture(
            io.BytesIO(original_bytes), Inches(1), Inches(1), Inches(2), Inches(2)
        )
        slide.shapes.add_picture(
            io.BytesIO(make_test_image_bytes()), pic.left, pic.top, pic.width, pic.height
        )
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        prs.save(path)
        try:
            r = parse_file(path, PPTX_MIME, "occ3", "test.pptx", "qualitative")
            assert len(r["images"]) == 2
            assert r["images"][0]["image_bytes"] == original_bytes
        finally:
            os.unlink(path)

    def test_covered_by_earlier_shape_not_masked(self):
        # 덮는 도형이 그림보다 먼저(z-order상 아래) 그려졌다면 실제로는 그림이 위에 보이므로
        # 마스킹하면 안 된다.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cover = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2)
        )
        cover.fill.solid()
        original_bytes = make_test_image_bytes()
        slide.shapes.add_picture(
            io.BytesIO(original_bytes), cover.left, cover.top, cover.width, cover.height
        )
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        prs.save(path)
        try:
            r = parse_file(path, PPTX_MIME, "occ4", "test.pptx", "qualitative")
            assert len(r["images"]) == 1
            assert r["images"][0]["image_bytes"] == original_bytes
        finally:
            os.unlink(path)

    def test_covered_by_shape_in_outer_group_masked(self):
        # 실제 사례: 그림은 안쪽 그룹에 있고, 그 그림을 가리는 단색 도형은 바깥쪽 그룹의
        # 형제 도형으로 존재 (그림의 "같은 컨테이너"가 아님). 그룹 계층을 넘어서도 겹침을
        # 감지해야 한다.
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        original_bytes = make_test_image_bytes((60, 60))
        pic = slide.shapes.add_picture(
            io.BytesIO(original_bytes), Inches(1), Inches(1), Inches(2), Inches(2)
        )
        inner_group = slide.shapes.add_group_shape([pic])
        cover = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, pic.left, pic.top, Emu(int(pic.width / 2)), pic.height
        )
        cover.fill.solid()
        # 바깥쪽 그룹: [inner_group(그림 포함), cover] 순서 — cover가 나중에(위에) 그려짐
        slide.shapes.add_group_shape([inner_group, cover])
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        prs.save(path)
        try:
            r = parse_file(path, PPTX_MIME, "occ5", "test.pptx", "qualitative")
            assert len(r["images"]) == 1
            masked_bytes = r["images"][0]["image_bytes"]
            assert masked_bytes != original_bytes
            assert self._pixel(masked_bytes, (10, 30)) == (255, 255, 255)
            assert self._pixel(masked_bytes, (50, 30)) == self._pixel(original_bytes, (50, 30))
        finally:
            os.unlink(path)


# ── DOCX 테스트 ──────────────────────────────────────────


class TestDocxParser:
    def test_paragraphs_grouped_into_pages(self):
        # 100개 단락 → 페이지 2개 이상
        path = make_docx([f"단락 {i}" for i in range(100)])
        try:
            r = parse_file(path, DOCX_MIME, "f5", "test.docx", "qualitative")
            assert r["parse_error"] is None
            assert r["total_pages"] >= 2
            assert r["pages"][0]["page_number"] == 1
            assert r["pages"][1]["page_number"] == 2
        finally:
            os.unlink(path)

    def test_small_doc_single_page(self):
        path = make_docx([f"단락 {i}" for i in range(10)])
        try:
            r = parse_file(path, DOCX_MIME, "f6", "test.docx", "qualitative")
            assert r["total_pages"] == 1
        finally:
            os.unlink(path)

    def test_table_cells_extracted(self):
        path = make_docx(["본문 단락"], table_cells=["셀A", "셀B", "셀C"])
        try:
            r = parse_file(path, DOCX_MIME, "f7", "table.docx", "qualitative")
            combined = "\n".join(p["text"] for p in r["pages"])
            assert "셀A" in combined
            assert "셀B" in combined
            assert "셀C" in combined
        finally:
            os.unlink(path)

    def test_page_number_one_based(self):
        path = make_docx([f"단락 {i}" for i in range(60)])
        try:
            r = parse_file(path, DOCX_MIME, "f8", "test.docx", "qualitative")
            assert r["pages"][0]["page_number"] == 1
        finally:
            os.unlink(path)


# ── PDF 테스트 ──────────────────────────────────────────


class TestPdfParser:
    def test_page_count(self):
        path = make_pdf(["Page one", "Page two", "Page three"])
        try:
            r = parse_file(path, PDF_MIME, "f9", "test.pdf", "qualitative")
            assert r["parse_error"] is None
            assert r["total_pages"] == 3
        finally:
            os.unlink(path)

    def test_page_numbers_one_based(self):
        path = make_pdf(["First page"])
        try:
            r = parse_file(path, PDF_MIME, "f10", "test.pdf", "qualitative")
            assert r["pages"][0]["page_number"] == 1
        finally:
            os.unlink(path)

    def test_text_extracted(self):
        path = make_pdf(["Hello World"])
        try:
            r = parse_file(path, PDF_MIME, "f11", "test.pdf", "qualitative")
            assert r["pages"][0]["text"] != ""
        finally:
            os.unlink(path)

    def test_scanned_pdf_returns_error(self):
        # 빈 페이지만 5개 → 100% 비어있음 → 스캔 PDF 오류
        path = make_pdf(["", "", "", "", ""])
        try:
            r = parse_file(path, PDF_MIME, "f12", "scanned.pdf", "qualitative")
            assert r["parse_error"] is not None
            assert "텍스트 레이어" in r["parse_error"]
            assert r["pages"] == []
            assert r["total_pages"] == 0
        finally:
            os.unlink(path)

    def test_mostly_empty_pdf_returns_error(self):
        # 페이지 5개 중 4개 비어있음 (80%) → 오류
        path = make_pdf(["text", "", "", "", ""])
        try:
            r = parse_file(path, PDF_MIME, "f13", "scan2.pdf", "qualitative")
            assert r["parse_error"] is not None
        finally:
            os.unlink(path)

    def test_mixed_pdf_ok(self):
        # 페이지 5개 중 1개만 비어있음 (20%) → 정상
        path = make_pdf(["p1", "p2", "p3", "p4", ""])
        try:
            r = parse_file(path, PDF_MIME, "f14", "mixed.pdf", "qualitative")
            assert r["parse_error"] is None
            assert r["total_pages"] == 5
        finally:
            os.unlink(path)


# ── 지원하지 않는 형식 ────────────────────────────────────


class TestUnsupportedFormats:
    def test_ppt_error(self):
        r = parse_file("dummy.ppt", "application/vnd.ms-powerpoint", "f15", "old.ppt", "qualitative")
        assert r["parse_error"] is not None
        assert "구버전 PPT" in r["parse_error"]
        assert r["pages"] == []

    def test_doc_error(self):
        r = parse_file("dummy.doc", "application/msword", "f16", "old.doc", "qualitative")
        assert r["parse_error"] is not None
        assert "구버전 DOC" in r["parse_error"]
        assert r["pages"] == []

    def test_unknown_extension_error(self):
        r = parse_file("dummy.xyz", "application/octet-stream", "f17", "file.xyz", "qualitative")
        assert r["parse_error"] is not None
