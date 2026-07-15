import io
import os
from typing import Optional

PARAS_PER_PAGE = 50

_EXT_UNSUPPORTED = {
    ".ppt": "구버전 PPT 형식은 지원하지 않습니다. PPTX로 변환 후 업로드하세요.",
    ".doc": "구버전 DOC 형식은 지원하지 않습니다. DOCX로 변환 후 업로드하세요.",
}

_MIME_EXT = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.ms-powerpoint": ".ppt",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/msword": ".doc",
    "application/pdf": ".pdf",
}

MIN_IMAGE_BYTES = 1_000  # 1 KB 미만 아이콘·불릿 이미지 제외
MAX_IMAGES_PER_FILE = 40  # 파일당 최대 처리 이미지 수 (신규 콘텐츠/마스터 이미지 종류당 1회분 예산)
MAX_MASTER_LOGO_REPEATS = 60  # 이미 감지된 마스터/레이아웃 로고가 슬라이드마다 반복 등록될 때의 별도 예산
MAX_PDF_PAGE_RENDERS = 30  # PDF 페이지 렌더 최대 수 (벡터 로고 보완용)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _parse_pptx(file_path: str) -> tuple[list[dict], list[dict], Optional[str]]:
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def _shape_bbox(shape) -> Optional[tuple[int, int, int, int]]:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        if None in (left, top, width, height) or width <= 0 or height <= 0:
            return None
        return left, top, width, height

    def _intersect(
        a: tuple[int, int, int, int], b: tuple[int, int, int, int]
    ) -> Optional[tuple[int, int, int, int]]:
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        x1, y1 = max(ax, bx), max(ay, by)
        x2, y2 = min(ax + aw, bx + bw), min(ay + ah, by + bh)
        if x2 <= x1 or y2 <= y1:
            return None
        return x1, y1, x2 - x1, y2 - y1

    def _is_opaque_cover(shape) -> bool:
        # 그림(Picture)은 가운데가 뚫린 액자·오버레이일 수 있어 픽셀 확인 없이는 불투명 여부를
        # 신뢰할 수 없으므로 커버로 취급하지 않는다. 단색 채우기 도형(사각형 등으로 가린 경우)만
        # 안전하게 "가림"으로 판단한다.
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return False
        try:
            return shape.fill.type == MSO_FILL_TYPE.SOLID
        except Exception:
            return False

    def _apply_transform_stack(
        bbox: tuple[int, int, int, int],
        stack: list[tuple[int, int, int, int, int, int, int, int]],
    ) -> tuple[int, int, int, int]:
        # 그룹 안에 중첩된 도형의 left/top/width/height는 python-pptx가 슬라이드 절대좌표가
        # 아니라 그 그룹의 chOff/chExt 기준 로컬 좌표를 그대로 돌려준다 (그룹이 배율 조정되어
        # 있으면 실제 위치·크기와 달라짐). 바깥쪽 그룹부터 순서대로 off/ext ↔ chOff/chExt
        # 변환을 누적 적용해 슬라이드 절대좌표로 환산해야 서로 다른 그룹에 속한 도형끼리도
        # (예: 인증서 그림과 그 위를 덮는 사각형이 다른 그룹 계층에 있는 경우) 정확히 겹침을
        # 판정할 수 있다.
        x, y, w, h = bbox
        x2, y2 = x + w, y + h
        for off_x, off_y, ext_cx, ext_cy, choff_x, choff_y, chext_cx, chext_cy in reversed(stack):
            sx = ext_cx / chext_cx if chext_cx else 1
            sy = ext_cy / chext_cy if chext_cy else 1
            x = off_x + (x - choff_x) * sx
            y = off_y + (y - choff_y) * sy
            x2 = off_x + (x2 - choff_x) * sx
            y2 = off_y + (y2 - choff_y) * sy
        return x, y, x2 - x, y2 - y

    def _flatten_shapes(shapes, stack: list) -> list[tuple]:
        """슬라이드(또는 레이아웃/마스터) 하나의 모든 도형을 전역 z-order(문서 순서) 그대로
        평탄화하면서 각 도형의 절대좌표 bbox를 계산한다.

        그룹은 그 자체로는 그려지는 대상이 아니라 컨테이너이므로 결과에 포함하지 않고,
        하위 도형만 (도형, 절대bbox) 튜플로 담는다. 그룹의 자식은 부모의 z-order 안에
        그 그룹이 위치한 자리에 그대로 끼워지므로, 문서 순서대로 재귀 평탄화한 목록이 곧
        전역 페인트 순서와 일치한다.
        """
        flat: list[tuple] = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                xfrm = shape._element.xfrm
                if xfrm is None:
                    flat.extend(_flatten_shapes(shape.shapes, stack))
                    continue
                entry = (
                    shape.left, shape.top, shape.width, shape.height,
                    shape._element.chOff.x, shape._element.chOff.y,
                    shape._element.chExt.cx, shape._element.chExt.cy,
                )
                flat.extend(_flatten_shapes(shape.shapes, stack + [entry]))
            else:
                bbox = _shape_bbox(shape)
                abs_bbox = _apply_transform_stack(bbox, stack) if bbox is not None else None
                flat.append((shape, abs_bbox))
        return flat

    def _covering_rects(flat: list[tuple], index: int, bbox: tuple[int, int, int, int]) -> list[tuple[int, int, int, int]]:
        """이 그림보다 나중에 그려진(z-order 상 위) 불투명 도형들이 이 그림과 겹치는
        절대좌표 영역을 모두 모은다.

        문서형 이미지(인증서 등)는 회사 식별정보가 이미지 전체가 아니라 일부 문단에만 있는
        경우가 많아, "이미지 전체가 거의 다 가려졌는지"만으로는 부분 가림(레닥션)을 못 잡는다.
        대신 겹치는 영역 좌표를 반환해 호출자가 그 부분만 실제로 마스킹하도록 한다.
        """
        rects = []
        for later_shape, later_bbox in flat[index + 1:]:
            if later_bbox is None or not _is_opaque_cover(later_shape):
                continue
            inter = _intersect(bbox, later_bbox)
            if inter is not None:
                rects.append(inter)
        return rects

    def _mask_covered_regions(
        blob: bytes, bbox: tuple[int, int, int, int], covers: list[tuple[int, int, int, int]]
    ) -> Optional[bytes]:
        """그림 위에 얹힌 불투명 도형에 가려지는 영역을 실제로 흰색으로 칠해, 화면에 실제
        보이는 모습 그대로 AI 검출에 넘긴다. 디코딩 실패 시 None을 반환해 원본을 그대로 쓰게 한다."""
        try:
            from PIL import Image, ImageDraw
            img = Image.open(io.BytesIO(blob)).convert("RGB")
        except Exception:
            return None
        w, h = img.size
        bx, by, bw, bh = bbox
        draw = ImageDraw.Draw(img)
        for cx, cy, cw, ch in covers:
            left = round((cx - bx) / bw * w)
            top = round((cy - by) / bh * h)
            right = round((cx + cw - bx) / bw * w)
            bottom = round((cy + ch - by) / bh * h)
            draw.rectangle([left, top, right, bottom], fill=(255, 255, 255))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()

    def _extract_shape_texts(shape) -> list[str]:
        texts = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                texts.extend(_extract_shape_texts(s))
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        if para.text:
                            texts.append(para.text)
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text:
                    texts.append(para.text)
        return texts

    def _extract_container_images(shapes, page_num: int) -> list[dict]:
        """도형 트리(슬라이드 최상위, 레이아웃, 마스터 중 하나) 전체에서 그림을 추출.

        전역 z-order(문서 순서) 기준으로, 나중에 그려진 불투명 도형에 가려지는 영역이 있으면
        그 부분을 흰색으로 마스킹한 뒤 추출한다 — 그룹 계층이 달라도(예: 그림은 안쪽 그룹,
        가리는 도형은 바깥쪽 그룹) 정확히 감지되도록 슬라이드 절대좌표로 비교한다.
        """
        flat = _flatten_shapes(shapes, [])
        imgs = []
        for i, (shape, abs_bbox) in enumerate(flat):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                mime_type = shape.image.content_type or "image/png"
                if len(blob) < MIN_IMAGE_BYTES:
                    continue
                if abs_bbox is not None:
                    covers = _covering_rects(flat, i, abs_bbox)
                    if covers:
                        masked = _mask_covered_regions(blob, abs_bbox, covers)
                        if masked is not None:
                            blob, mime_type = masked, "image/png"
                imgs.append({
                    "page_number": page_num,
                    "image_bytes": blob,
                    "mime_type": mime_type,
                })
            except Exception:
                pass
        return imgs

    prs = Presentation(file_path)
    pages, images = [], []
    seen: set[tuple] = set()          # (img_hash, page_num) — 슬라이드별 중복 추가 방지
    seen_hashes: set[int] = set()      # img_hash — 콘텐츠/마스터 통틀어 이미 예산을 소모한 이미지인지 판단
    master_repeat_count = 0

    def _add_shapes(shapes, page_num: int) -> None:
        for img in _extract_container_images(shapes, page_num):
            if len(images) >= MAX_IMAGES_PER_FILE:
                return
            h = hash(img["image_bytes"])
            key = (h, page_num)
            if key not in seen:
                seen.add(key)
                seen_hashes.add(h)
                images.append(img)

    def _add_master_shapes(shapes, page_num: int) -> None:
        # 마스터/레이아웃 이미지도 슬라이드별로 추가 — 같은 로고라도 슬라이드마다 page_num이 달라야
        # 올바른 페이지에 검출 결과가 기록된다.
        # 단, 같은 마스터 로고가 슬라이드 수만큼 반복되면 콘텐츠 이미지 예산(MAX_IMAGES_PER_FILE)을
        # 잠식하므로, 최초 발견 시에만 콘텐츠 예산을 사용하고 이후 반복 등장은 별도의
        # MAX_MASTER_LOGO_REPEATS 예산으로 관리한다.
        nonlocal master_repeat_count
        for img in _extract_container_images(shapes, page_num):
            h = hash(img["image_bytes"])
            key = (h, page_num)
            if key in seen:
                continue
            if h in seen_hashes:
                if master_repeat_count >= MAX_MASTER_LOGO_REPEATS:
                    continue
                seen.add(key)
                images.append(img)
                master_repeat_count += 1
            else:
                if len(images) >= MAX_IMAGES_PER_FILE:
                    continue
                seen.add(key)
                seen_hashes.add(h)
                images.append(img)

    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            texts.extend(_extract_shape_texts(shape))
        pages.append({"page_number": idx, "text": "\n".join(texts)})

        _add_shapes(slide.shapes, idx)
        # 레이아웃·마스터에 있는 이미지도 각 슬라이드에 포함 (공통 로고 등)
        _add_master_shapes(slide.slide_layout.shapes, idx)
        _add_master_shapes(slide.slide_layout.slide_master.shapes, idx)

    return pages, images, None


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _parse_docx(file_path: str) -> tuple[list[dict], list[dict], Optional[str]]:
    from docx import Document

    doc = Document(file_path)
    paragraphs: list[str] = []
    for para in doc.paragraphs:
        paragraphs.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    paragraphs.append(para.text)

    if not paragraphs:
        return [{"page_number": 1, "text": ""}], [], None

    pages = []
    for i in range(0, len(paragraphs), PARAS_PER_PAGE):
        chunk = paragraphs[i: i + PARAS_PER_PAGE]
        pages.append({
            "page_number": i // PARAS_PER_PAGE + 1,
            "text": "\n".join(chunk),
        })

    images = []
    try:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype and len(images) < MAX_IMAGES_PER_FILE:
                blob = rel.target_part.blob
                if len(blob) >= MIN_IMAGE_BYTES:
                    images.append({
                        "page_number": 1,
                        "image_bytes": blob,
                        "mime_type": rel.target_part.content_type or "image/png",
                    })
    except Exception:
        pass

    return pages, images, None


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_pdf(file_path: str) -> tuple[list[dict], list[dict], Optional[str]]:
    import pdfplumber

    pages = []
    with pdfplumber.open(file_path) as pdf:
        total = len(pdf.pages)
        if total == 0:
            return [], [], None

        empty_count = 0
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if not text.strip():
                empty_count += 1
            pages.append({"page_number": idx, "text": text})

        if empty_count / total >= 0.8:
            non_empty = [p for p in pages if p["text"].strip()]
            if not non_empty:
                return [], [], "텍스트 레이어를 찾을 수 없습니다. 스캔된 PDF일 수 있습니다."
            return non_empty, _extract_pdf_images(file_path), "일부 페이지에 텍스트가 없습니다. 스캔 페이지가 포함되어 있을 수 있습니다."

    images = _extract_pdf_images(file_path)
    return pages, images, None


def _extract_pdf_images(file_path: str) -> list[dict]:
    try:
        import fitz  # pymupdf
    except ImportError:
        return []

    images = []
    try:
        doc = fitz.open(file_path)

        # 1. 임베드된 래스터 이미지 추출
        for page_idx, page in enumerate(doc, start=1):
            for img_info in page.get_images(full=False):
                if len(images) >= MAX_IMAGES_PER_FILE:
                    break
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    blob = base_image["image"]
                    if len(blob) >= MIN_IMAGE_BYTES:
                        ext = base_image.get("ext", "png")
                        images.append({
                            "page_number": page_idx,
                            "image_bytes": blob,
                            "mime_type": f"image/{ext}",
                        })
                except Exception:
                    continue

        # 2. 페이지 렌더링 추가 — 벡터 로고·Form XObject 내 이미지를 보완
        #    임베드 추출로 감지되지 않는 로고를 페이지 전체 이미지로 재확인.
        render_count = 0
        for page_idx, page in enumerate(doc, start=1):
            if render_count >= MAX_PDF_PAGE_RENDERS:
                break
            try:
                mat = fitz.Matrix(1.0, 1.0)  # 72 DPI — 로고 식별에 충분한 해상도
                pix = page.get_pixmap(matrix=mat, alpha=False)
                blob = pix.tobytes("png")
                if blob:
                    images.append({
                        "page_number": page_idx,
                        "image_bytes": blob,
                        "mime_type": "image/png",
                        "is_page_render": True,
                    })
                    render_count += 1
            except Exception:
                continue
    except Exception:
        pass
    return images


# ── 진입점 ─────────────────────────────────────────────────────────────────────

def parse_file(
    file_path: str,
    mime_type: str,
    file_id: str,
    original_filename: str,
    proposal_type: str,
) -> dict:
    result: dict = {
        "file_id": file_id,
        "original_filename": original_filename,
        "proposal_type": proposal_type,
        "pages": [],
        "images": [],
        "total_pages": 0,
        "parse_error": None,
    }

    ext = os.path.splitext(original_filename)[1].lower()
    if not ext:
        ext = _MIME_EXT.get(mime_type, "")

    if ext in _EXT_UNSUPPORTED:
        result["parse_error"] = _EXT_UNSUPPORTED[ext]
        return result

    try:
        if ext == ".pptx":
            pages, images, error = _parse_pptx(file_path)
        elif ext == ".docx":
            pages, images, error = _parse_docx(file_path)
        elif ext == ".pdf":
            pages, images, error = _parse_pdf(file_path)
        else:
            result["parse_error"] = f"지원하지 않는 파일 형식입니다: {ext or mime_type}"
            return result

        result["pages"] = pages
        result["images"] = images
        result["total_pages"] = len(pages)
        result["parse_error"] = error
    except Exception as e:
        result["parse_error"] = f"파일 파싱 중 오류가 발생했습니다: {str(e)}"

    return result
